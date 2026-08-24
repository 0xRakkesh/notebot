import os
import re
import sys
import io
import time
import requests
from typing import TypedDict, List, Dict
from tavily import TavilyClient

from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from langgraph.graph import StateGraph, START, END
from langgraph.checkpoint.memory import MemorySaver

# Force UTF-8 for Windows console
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
load_dotenv()

# ==========================================
# 1. State Definition
# ==========================================
class AgentState(TypedDict):
    youtube_url: str
    file_path: str
    transcript: str
    error: str
    key_concepts: List[str]
    diagrams: Dict[str, str]
    final_notes: str

# ==========================================
# 2. Setup LLMs & Tools
# ==========================================
# Using tokenin API
llm = ChatOpenAI(
    model="myt/claude-opus-4-8-free",
    api_key=os.getenv("TOKENIN_API_KEY"),
    base_url="https://tokenin.my.id/v1"
)
memory = MemorySaver()

# Global rate limit state
last_llm_call = 0.0

def wait_for_rate_limit():
    global last_llm_call
    now = time.time()
    elapsed = now - last_llm_call
    # We need 60 seconds between calls (for 1 req/min limit), plus 2s buffer
    if elapsed < 62.0 and last_llm_call != 0.0:
        wait_time = 62.0 - elapsed
        print(f"   [Rate Limit] Sleeping for {wait_time:.1f}s to respect API limits...")
        time.sleep(wait_time)
    
    last_llm_call = time.time()

def search_tavily_diagram(query: str) -> str:
    """Helper function to search Tavily for a diagram."""
    tavily_api_key = os.getenv("TAVILY_API_KEY")
    if not tavily_api_key:
        return "Error: TAVILY_API_KEY not found."
    client = TavilyClient(api_key=tavily_api_key)
    try:
        # Search specifically for diagrams, architecture, or flowcharts
        search_query = f"{query} architecture diagram OR flowchart OR architecture"
        response = client.search(search_query, search_depth="basic", include_images=True)
        if "images" in response and len(response["images"]) > 0:
            return response["images"][0]
        return ""
    except Exception as e:
        return ""

def fetch_youtube_subtitles_raw(video_url: str) -> str:
    """Original reliable fetcher using yt-dlp to extract URL, then requests.get"""
    try:
        import yt_dlp
        match = re.search(r"(?:v=|\/)([0-9A-Za-z_-]{11}).*", video_url)
        video_id = match.group(1) if match else video_url
        url = f"https://www.youtube.com/watch?v={video_id}"

        ydl_opts = {
            'skip_download': True,
            'writesubtitles': True,
            'writeautomaticsub': True,
            'subtitleslangs': ['en'],
            'quiet': True,
            'no_warnings': True,
            'ignore_no_formats_error': True,
            'format': 'worst*',
            'extractor_args': {'youtube': {'player_client': ['android', 'ios', 'web']}},
        }
        
        proxy = os.getenv("YOUTUBE_PROXY")
        if proxy:
            ydl_opts['proxy'] = proxy

        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
            duration = info.get('duration', 0)
            if duration > 3600:
                return f"Error: The video is longer than 1 hour ({int(duration/60)} mins). Rejecting the request."

            subs = info.get('subtitles', {})
            auto_subs = info.get('automatic_captions', {})

            sub_url = None
            for lang in ['en', 'en-US', 'en-GB', 'en-IN']:
                for source in [subs, auto_subs]:
                    if lang in source:
                        for fmt in source[lang]:
                            if fmt.get('ext') == 'json3':
                                sub_url = fmt['url']
                                break
                        if sub_url:
                            break
                if sub_url:
                    break

            if not sub_url:
                return "Error: No English transcript found for this video."

            resp = requests.get(sub_url)
            caption_json = resp.json()

            formatted = []
            for event in caption_json.get('events', []):
                start_ms = event.get('tStartMs', 0)
                segs = event.get('segs', [])
                if not segs:
                    continue
                text = ''.join(s.get('utf8', '') for s in segs).strip()
                if not text:
                    continue
                mins = int((start_ms / 1000) // 60)
                secs = int((start_ms / 1000) % 60)
                formatted.append(f"[{mins:02d}:{secs:02d}] {text}")

            if not formatted:
                return "Error: Transcript was empty."

            transcript_text = "\n".join(formatted)
            MAX_CHARS = 120000 
            if len(transcript_text) > MAX_CHARS:
                transcript_text = transcript_text[:MAX_CHARS] + "\n\n...[Transcript truncated due to length limits]..."
                
            return transcript_text
    except Exception as e:
        return f"Error fetching transcript: {str(e)}"

def extract_pdf_text(file_path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text.append(t)
        return "\n".join(text)
    except Exception as e:
        return f"Error: Failed to parse PDF - {e}"

def extract_pptx_text(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text.append("\n".join(slide_text))
        return "\n\n".join(text)
    except Exception as e:
        return f"Error: Failed to parse PPTX - {e}"

# ==========================================
# 3. Graph Nodes
# ==========================================
def node_fetch_content(state: AgentState) -> dict:
    if state.get("file_path"):
        print(f"-> Extracting document: {state['file_path']}")
        file_path = state["file_path"]
        if file_path.lower().endswith('.pdf'):
            content = extract_pdf_text(file_path)
        elif file_path.lower().endswith(('.pptx', '.ppt')):
            content = extract_pptx_text(file_path)
        else:
            return {"error": "Error: Unsupported file format."}
            
        if content.startswith("Error:"):
            return {"error": content}
        
        # Limit content size 
        if len(content) > 120000:
            content = content[:120000] + "\n\n...[Truncated due to limits]..."
            
        return {"transcript": content, "error": ""}
    elif state.get("youtube_url"):
        print("-> Fetching transcript...")
        transcript = fetch_youtube_subtitles_raw(state["youtube_url"])
        if transcript.startswith("Error:"):
            return {"error": transcript}
        return {"transcript": transcript, "error": ""}
    else:
        return {"error": "Error: No input provided."}

def node_research_expansion(state: AgentState) -> dict:
    if state.get("error"): return {}
    
    transcript = state.get("transcript", "")
    
    # Expand if it's a short document
    if state.get("file_path") and len(transcript) < 3000 and len(transcript) > 50:
        print("-> Document is very brief. Expanding via research...")
        
        prompt = f"""The following is text extracted from a short presentation or document. It likely contains high-level topics or bullet points. 
Generate exactly ONE comprehensive web search query that would yield detailed educational information about the core topics mentioned here. 
Return ONLY the search query.
Document: {transcript}"""
        
        try:
            wait_for_rate_limit()
            response = llm.invoke([HumanMessage(content=prompt)])
            query = response.content.strip().strip('"\'')
            print(f"   Generated search query: {query}")
            
            tavily_api_key = os.getenv("TAVILY_API_KEY")
            if tavily_api_key:
                client = TavilyClient(api_key=tavily_api_key)
                search_response = client.search(query, search_depth="advanced", max_results=2)
                
                research_text = "Additional Research Context:\n"
                for result in search_response.get("results", []):
                    research_text += f"{result.get('content', '')}\n\n"
                    
                print("   Appended research context to the document text.")
                return {"transcript": transcript + "\n\n" + research_text}
                
        except Exception as e:
            print(f"   Research failed: {e}")
            
    return {}

def node_extract_concepts(state: AgentState) -> dict:
    if state.get("error"): return {}
    
    print("-> Extracting key concepts...")
    # Truncate transcript to first 120,000 characters for concept extraction to save tokens
    extraction_transcript = state['transcript'][:120000]
    
    prompt = f"""Analyze the following video transcript. Identify the 1 to 3 MOST IMPORTANT highly-technical concepts, frameworks, or architectures discussed that would benefit from having a diagram or flowchart.
    
    Output ONLY a comma-separated list of these 1-3 concepts. Do not output anything else.
    
    Transcript:
    {extraction_transcript}
    """
    
    try:
        wait_for_rate_limit()
        response = llm.invoke([HumanMessage(content=prompt)])
        # Clean the output
        content = response.content.strip()
        # Split by comma
        concepts = [c.strip() for c in content.split(",") if c.strip()]
        # Take max 3
        concepts = concepts[:3]
        print(f"   Concepts found: {concepts}")
        return {"key_concepts": concepts}
    except Exception as e:
        print(f"   Failed to extract concepts: {e}")
        return {"key_concepts": []}

def node_fetch_diagrams(state: AgentState) -> dict:
    if state.get("error"): return {}
    
    print("-> Fetching diagrams...")
    diagrams = {}
    concepts = state.get("key_concepts", [])
    
    for concept in concepts:
        url = search_tavily_diagram(concept)
        if url and not url.startswith("Error"):
            diagrams[concept] = url
            print(f"   Found diagram for '{concept}': {url}")
            
    return {"diagrams": diagrams}

def node_generate_notes(state: AgentState) -> dict:
    if state.get("error"): return {}
    
    print("-> Generating final notes...")
    
    # Format the diagrams context
    diagrams_context = ""
    if state.get("diagrams"):
        diagrams_context = "Available Diagrams (Use these URLs when formatting the Architecture & Flow sections):\n"
        for concept, url in state["diagrams"].items():
            diagrams_context += f"- {concept}: {url}\n"
    else:
        diagrams_context = "No diagrams available."

    sys_prompt = """You are Marin Kitagawa, an anime CS nerd who makes the best notes because you are the topper of your class from a parallel anime universe. You extract core concepts from video transcripts into ultra-concise, Telegram-ready notes.
Topics will include OS, DBMS, Machine Learning, Image Processing, and Web Tech.
Goal: Create dense, skimmable cheat sheets focused on theory, architecture, code, and examples. Maintain your anime CS nerd persona slightly in the phrasing, but prioritize extreme technical density and accuracy.

RULES:
1. STRUCTURE: Do NOT use repetitive bullet points like "• Definition:". Use this clean format (leave a blank line between concepts):
   
   <b>CONCEPT NAME</b>
   <i>1-2 ultra-short sentences defining the concept.</i>
   
   <b>Example:</b> A minimal, practical real-world use case.
   
   <b>Architecture & Flow:</b>
   [If a diagram URL was provided for this concept in the 'Available Diagrams', output it EXACTLY like this: [IMAGE]url_from_tool[/IMAGE]. Otherwise describe the flow concisely.]
   
   <b>Code / Math:</b> (ONLY IF APPLICABLE, otherwise completely omit this section and do not output empty <pre> tags)
   <pre>
   [If the video shows code or formulas, provide the EXACT code/formula here. Do NOT include any code comments (no //, #, or /*).]
   </pre>
   <i>1-line explanations for key code parts. Use colons (:) to link explanations.</i>
2. CONCISENESS & NO BULLETS: Absolutely NO bullet points (no `•`, no `-`, no `*`) and NO arrows (no `→`, no `->`) anywhere. Use commas (,) and colons (:) to logically connect ideas, properties, or flow. Use newlines to separate larger points.
3. HIGHLIGHTING (STRICT HTML): You MUST use `<b>` tags for bold. NEVER use Markdown `**`. Example: <b>IMPORTANT</b>, not **IMPORTANT**.
4. BALANCE: Capture critical theory (like DBMS normalization, OS deadlocks, ML math) just as heavily as code.
5. REVISION: End with a <b>REVISION:</b> section containing 3-5 crucial technical takeaways (separated by newlines, NO bullets).
"""

    chunk_size = 120000
    transcript = state['transcript']
    chunks = [transcript[i:i+chunk_size] for i in range(0, len(transcript), chunk_size)]
    
    final_notes_pieces = []
    
    try:
        for i, chunk in enumerate(chunks):
            print(f"   Processing chunk {i+1}/{len(chunks)}...")
            
            # Remove REVISION section requirement for intermediate chunks to prevent duplicates
            chunk_sys_prompt = sys_prompt
            if len(chunks) > 1 and i < len(chunks) - 1:
                chunk_sys_prompt = sys_prompt.replace(
                    "5. REVISION: End with a <b>REVISION:</b> section containing 3-5 crucial technical takeaways (separated by newlines, NO bullets).",
                    "5. Do NOT output a REVISION section for this chunk."
                )
                
            human_prompt = f"{diagrams_context}\n\nTranscript (Part {i+1}/{len(chunks)}):\n{chunk}"
            
            wait_for_rate_limit()
            response = llm.invoke([
                SystemMessage(content=chunk_sys_prompt),
                HumanMessage(content=human_prompt)
            ])
            
            final_notes_pieces.append(response.content)
            
        return {"final_notes": "\n\n---\n\n".join(final_notes_pieces)}
    except Exception as e:
        return {"error": f"Error generating notes: {e}"}

# ==========================================
# 4. Build Graph
# ==========================================
workflow = StateGraph(AgentState)

workflow.add_node("fetch", node_fetch_content)
workflow.add_node("research", node_research_expansion)
workflow.add_node("extract", node_extract_concepts)
workflow.add_node("diagrams", node_fetch_diagrams)
workflow.add_node("generate", node_generate_notes)

workflow.add_edge(START, "fetch")
workflow.add_edge("fetch", "research")
workflow.add_edge("research", "extract")
workflow.add_edge("extract", "diagrams")
workflow.add_edge("diagrams", "generate")
workflow.add_edge("generate", END)

advanced_agent = workflow.compile(checkpointer=memory)

# ==========================================
# 5. Interface
# ==========================================
def process_video(youtube_url: str, session_id: str = "notebot_conversation_1") -> str:
    """Processes a video URL through the advanced graph and returns the generated notes."""
    config = {"configurable": {"thread_id": session_id}}
    
    # Initialize the state
    inputs = {
        "youtube_url": youtube_url,
        "file_path": "",
        "transcript": "",
        "error": "",
        "key_concepts": [],
        "diagrams": {},
        "final_notes": ""
    }
    
    try:
        result = advanced_agent.invoke(inputs, config=config)
        if result.get("error"):
            return result["error"]
        
        return result.get("final_notes", "Error: No notes generated.")
    except Exception as e:
        return f"Error processing video: {e}"

def process_document(file_path: str, session_id: str = "notebot_conversation_1") -> str:
    """Processes a document (PDF/PPTX) through the advanced graph and returns the generated notes."""
    config = {"configurable": {"thread_id": session_id}}
    
    inputs = {
        "youtube_url": "",
        "file_path": file_path,
        "transcript": "",
        "error": "",
        "key_concepts": [],
        "diagrams": {},
        "final_notes": ""
    }
    
    try:
        result = advanced_agent.invoke(inputs, config=config)
        if result.get("error"):
            return result["error"]
        
        return result.get("final_notes", "Error: No notes generated.")
    except Exception as e:
        return f"Error processing document: {e}"

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: uv run src/agent/agent.py <youtube_url>")
        sys.exit(1)

    url = sys.argv[1]
    print(f"\n--- NoteBot: Advanced Agent Processing ---")
    print(f"URL: {url}\n")

    notes = process_video(url)
    print("\n--- Notes Generated ---\n")
    print(notes)
